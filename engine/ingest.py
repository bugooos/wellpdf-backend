# engine/ingest.py

from pdfminer.high_level import extract_text
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
import pandas as pd
import json
import math


def _safe(val):
    if val is None:
        return ""
    if isinstance(val, float) and math.isnan(val):
        return ""
    return str(val)


def ingest(path, ext):
    ext = ext.lower()

    # -------- TEXT --------
    if ext in ("txt", "log", "md"):
        return "text", open(path, encoding="utf-8", errors="ignore").read()

    if ext == "html":
        soup = BeautifulSoup(open(path, encoding="utf-8", errors="ignore"), "html.parser")
        return "text", soup.get_text()

    if ext == "pdf":
        return "text", extract_text(path)

    if ext in ("doc", "docx", "odt", "rtf"):
        doc = Document(path)
        return "text", "\n".join(p.text for p in doc.paragraphs)

    # -------- SPREADSHEETS (FIXED) --------
    if ext in ("csv", "xls", "xlsx", "ods"):
        if ext == "csv":
            df = pd.read_csv(path, dtype=str)
        else:
            df = pd.read_excel(path, dtype=str)

        df = df.fillna("")

        header = [_safe(c) for c in df.columns.tolist()]
        rows = [[_safe(v) for v in row] for row in df.values.tolist()]

        return "table", [header] + rows

    # -------- JSON --------
    if ext == "json":
        obj = json.load(open(path, encoding="utf-8", errors="ignore"))

        if isinstance(obj, list) and obj and isinstance(obj[0], dict):
            headers = list(obj[0].keys())
            rows = [[_safe(v) for v in o.values()] for o in obj]
            return "table", [headers] + rows

        return "text", json.dumps(obj, indent=2)

    # -------- XML --------
    if ext == "xml":
        soup = BeautifulSoup(open(path, encoding="utf-8", errors="ignore"), "xml")
        return "text", soup.get_text()

    # -------- PPT --------
    if ext in ("ppt", "pptx", "odp"):
        from pptx import Presentation
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "text", "\n".join(text)

    # -------- EPUB --------
    if ext == "epub":
        book = epub.read_epub(path)
        text = []
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text.append(soup.get_text())
        return "text", "\n".join(text)

    raise ValueError("Unsupported format")
