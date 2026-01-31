# engine/emit.py

from reportlab.lib.pagesizes import A4
from odf.opendocument import OpenDocumentText
from odf.text import P
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Table,
    TableStyle,
    Spacer,
)
from reportlab.lib import colors

from docx import Document
import pandas as pd
import json, csv
from io import StringIO
from ebooklib import epub


# =========================
# HELPERS
# =========================

def normalize_text(data):
    """Convert any data into safe text"""
    if isinstance(data, str):
        return data
    if isinstance(data, list):
        return "\n".join(" ".join(map(str, row)) for row in data)
    return str(data)


# =========================
# EMIT
# =========================

def emit(kind, data, out_path, ext):
    ext = ext.lower()

    # ========================
    # TEXT OUTPUTS
    # ========================
    if ext in ("txt", "log", "md"):
        text = normalize_text(data)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(text)
        return

    if ext == "html":
        text = normalize_text(data)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(f"<pre>{text}</pre>")
        return

    # ========================
    # PDF OUTPUT (PROPER)
    # ========================
    if ext == "pdf":
        doc = SimpleDocTemplate(
            out_path,
            pagesize=A4,
            leftMargin=40,
            rightMargin=40,
            topMargin=40,
            bottomMargin=40,
        )

        styles = getSampleStyleSheet()
        normal = ParagraphStyle(
            "NormalWrap",
            parent=styles["Normal"],
            fontSize=9,
            leading=12,
            wordWrap="LTR",
        )

        story = []

        # ---- TABLE → PDF ----
        if kind == "table" and isinstance(data, list):
            col_count = len(data[0])
            page_width = A4[0] - doc.leftMargin - doc.rightMargin
            col_widths = [page_width / col_count] * col_count

            table = Table(data, colWidths=col_widths, repeatRows=1)
            table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
            ]))

            story.append(table)

        # ---- TEXT → PDF ----
        else:
            text = normalize_text(data)
            for line in text.splitlines():
                story.append(Paragraph(line or "&nbsp;", normal))
                story.append(Spacer(1, 6))

        doc.build(story)
        return

    # ========================
    # WORD
    # ========================
    if ext in ("doc", "docx"):
        text = normalize_text(data)
        doc = Document()
        for line in text.splitlines():
            doc.add_paragraph(line)
        doc.save(out_path)
        return

    # ========================
    # CSV
    # ========================
    if ext == "csv":
        with open(out_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            if kind == "table" and isinstance(data, list):
                writer.writerows(data)
            else:
                for line in normalize_text(data).splitlines():
                    writer.writerow([line])
        return

    # ========================
    # SPREADSHEETS
    # ========================
    if ext in ("xls", "xlsx", "ods"):
        if kind == "table" and isinstance(data, list):
            df = pd.DataFrame(data[1:], columns=data[0])
        else:
            df = pd.read_csv(StringIO(normalize_text(data)))
        df.to_excel(out_path, index=False)
        return

    # ========================
    # JSON
    # ========================
    if ext == "json":
        if kind == "table" and isinstance(data, list):
            rows = [dict(zip(data[0], r)) for r in data[1:]]
            json.dump(rows, open(out_path, "w", encoding="utf-8"), indent=2)
        else:
            json.dump(
                {"content": normalize_text(data).splitlines()},
                open(out_path, "w", encoding="utf-8"),
                indent=2,
            )
        return

    # ========================
    # XML
    # ========================
    if ext == "xml":
        text = normalize_text(data)
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(
                "<document>\n" +
                "\n".join(f"<line>{line}</line>" for line in text.splitlines()) +
                "\n</document>"
            )
        return

    # ========================
    # PRESENTATION
    # ========================
    if ext in ("ppt", "pptx", "odp"):
        from pptx import Presentation
        text = normalize_text(data)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Converted Content"
        slide.placeholders[1].text = text[:4000]
        prs.save(out_path)
        return

    # ========================
    # EPUB
    # ========================
    if ext == "epub":
        text = normalize_text(data)

        book = epub.EpubBook()
        book.set_title("Converted File")
        content = epub.EpubHtml(title="Content", file_name="content.xhtml")
        content.content = f"<pre>{text}</pre>"

        book.add_item(content)
        book.spine = ["nav", content]
        epub.write_epub(out_path, book)
        return
    
        # ========================
    # ODT
    # ========================
    if ext == "odt":
        text = normalize_text(data)

        doc = OpenDocumentText()
        for line in text.splitlines():
            doc.text.addElement(P(text=line))

        doc.save(out_path)
        return
    
        # ========================
    # RTF
    # ========================
    if ext == "rtf":
        text = normalize_text(data)

        def rtf_escape(s):
            return s.replace("\\", r"\\").replace("{", r"\{").replace("}", r"\}")

        with open(out_path, "w", encoding="utf-8") as f:
            f.write(r"{\rtf1\ansi\deff0" + "\n")
            for line in text.splitlines():
                f.write(rtf_escape(line) + r"\par" + "\n")
            f.write("}")
        return

    raise ValueError("Emit unsupported")
