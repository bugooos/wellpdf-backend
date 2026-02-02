from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from pdfminer.high_level import extract_text
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import shutil, os, re

from core.utils import temp_input_path, safe_output_path

router = APIRouter()


@router.post("/convert/pdf-to-pptx", summary="Convert PDF to PPTX (formatted)")
async def pdf_to_pptx(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files allowed")

    input_path = temp_input_path("pdf")
    output_path = safe_output_path(file.filename, "pptx")

    with open(input_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    try:
        raw_text = extract_text(input_path)
        if not raw_text.strip():
            raise RuntimeError("No readable text found in PDF")

        prs = Presentation()

        SLIDE_WIDTH = prs.slide_width
        SLIDE_HEIGHT = prs.slide_height

        # Safe margins
        LEFT = Inches(1)
        TOP = Inches(1.2)
        WIDTH = SLIDE_WIDTH - Inches(2)
        HEIGHT = SLIDE_HEIGHT - Inches(2.2)

        MAX_LINES = 10

        lines = []
        for line in raw_text.splitlines():
            clean = line.strip()
            if clean:
                lines.append(clean)

        slides_content = []
        current = []

        for line in lines:
            current.append(line)
            if len(current) >= MAX_LINES:
                slides_content.append(current)
                current = []

        if current:
            slides_content.append(current)

        for idx, slide_lines in enumerate(slides_content):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide

            # Title
            title_box = slide.shapes.add_textbox(
                LEFT, Inches(0.4), WIDTH, Inches(0.6)
            )
            title_tf = title_box.text_frame
            title_tf.clear()

            title_p = title_tf.paragraphs[0]
            title_p.text = f"Slide {idx + 1}"
            title_p.font.size = Pt(24)
            title_p.font.bold = True
            title_p.alignment = PP_ALIGN.LEFT

            # Content box
            content_box = slide.shapes.add_textbox(
                LEFT, TOP, WIDTH, HEIGHT
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.clear()

            for i, line in enumerate(slide_lines):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = line
                p.font.size = Pt(16)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.LEFT

                # Bullet detection
                if re.match(r"^[-•]", line):
                    p.text = line.lstrip("-• ").strip()
                    p.level = 1
                else:
                    p.level = 0

        prs.save(output_path)

        return FileResponse(
            output_path,
            filename=os.path.basename(output_path),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    finally:
        if os.path.exists(input_path):
            os.remove(input_path)
