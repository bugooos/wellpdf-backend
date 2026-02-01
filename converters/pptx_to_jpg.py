from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from PIL import Image, ImageDraw
import shutil, os

from core.utils import temp_input_path, safe_output_path

router = APIRouter()

@router.post("/convert/pptx-to-jpg", summary="Convert PPTX to JPG (text)")
async def pptx_to_jpg(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only PPTX files allowed")

    input_path = temp_input_path("pptx")
    output_path = safe_output_path(file.filename, "jpg")

    with open(input_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    try:
        prs = Presentation(input_path)
        img = Image.new("RGB", (1200, 800), "white")
        draw = ImageDraw.Draw(img)

        y = 40
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    draw.text((40, y), shape.text, fill="black")
                    y += 20
            y += 40

        img.save(output_path, "JPEG", quality=95)
        return FileResponse(output_path, filename=os.path.basename(output_path))

    finally:
        if os.path.exists(input_path):
            os.remove(input_path)
